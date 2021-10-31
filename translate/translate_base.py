# coding: UTF-8
import re

import lxml.etree
import lxml.html

import html
import json
import os
import sys
import warnings
from collections import defaultdict
from docx import Document
from google.cloud.translate_v2 import Client
from openpyxl import load_workbook
from openpyxl.cell import Cell
from pptx import Presentation
from pptx.oxml import CT_TextLineBreak


CREDS = os.environ["GOOGLE_APPLICATION_CREDENTIALS"]
DOCX_STYLE_PROPERTY = {
    'font': {
        'color': {
            'rgb': {}, 'theme_color': {}, 'type': {},
        },
        'name': {}, 'size': {}, 'bold': {}, 'italic': {}, 'underline': {}, 'strike': {},
        'double_strike': {}, 'subscript': {}, 'superscript': {}, 'no_proof': {}, 'hidden': {},
        'highlight_color': {}, 'math': {}, 'rtl': {}, 'outline': {}, 'emboss': {}, 'shadow': {}, 'imprint': {},
        'all_caps': {}, 'small_caps': {}, 'web_hidden': {}, 'complex_script': {}, 'cs_bold': {}, 'cs_italic': {},
        'snap_to_grid': {}, 'spec_vanish': {},
    },
}
LANGUAGE_PROPERTIES = {
    'en': {'lang': 'English'},
    'fr': {'lang': 'French'},
    'ja': {
        'font': 'Hiragino Sans W3',
        # could also use 'MS Gothic', and for Windows could only find 'Hiragino Sans GB W3'
        'lang': 'Japanese',
        'overwrite': {
            ' / ': '/',  # what was this for?
        },
    }
}
TRANSLATION_RULES = {
    'ignore_styles': ['Code', 'Normal-No spellcheck', 'Legal Text'],
}


class GoogleTranslate(object):
    """ Establish a Google Cloud Translate client to translate passages of text. """
    def __init__(self, creds, target_lang, source_lang=None, online=True, history=None, show=0):
        self.client = Client.from_service_account_json(creds)
        self.target_lang = target_lang
        self.source_lang = source_lang
        self.online = online
        self.history = bool(history) and online
        self.history_file = history
        self.show = show
        self.prepare_translation()
        self.stats = defaultdict(dict)

    def prepare_translation(self):
        """ Reset stats; initialise translation dictionary using history if applicable and available. """
        self.cloud_requests = 0
        self.dummy_text = 0
        self.dict_hits = 0
        self.empty_strings = 0
        self.translated = {}
        if self.history:
            if os.path.exists(self.history_file):
                try:
                    with open(self.history_file, 'r') as f:
                        self.translated = json.load(f)
                except (IOError, OSError):
                    fallback_history_file = os.path.join(os.environ['HOME'], 'history_file.json')
                    warnings.warn("'%s' is a directory - saving history in home directory as:\n '%s' " % (
                        self.history_file, fallback_history_file
                    ))
                    self.history_file = fallback_history_file


    def request_translation(self, string):
        return self.client.translate(string, target_language=self.target_lang, source_language=self.source_lang)

    def translate(self, string):
        """ Translate a single text element. """
        if string and string not in self.translated:
            if self.show:
                print("{}{}".format(string[:self.show], " ..." if len(string) > self.show else ""))

            if self.online:
                translation = self.request_translation(string)['translatedText']
                # TODO: if not self.source_lang, look at ['detectedSourceLanguage']
                translation = html.unescape(translation)
                if self.target_lang in LANGUAGE_PROPERTIES and 'overwrite' in LANGUAGE_PROPERTIES[self.target_lang]:
                    for each in LANGUAGE_PROPERTIES[self.target_lang]['overwrite']:
                        if each in translation:
                            translation = LANGUAGE_PROPERTIES[self.target_lang]['overwrite'][each].join(
                                (translation.split(each))
                            )
                self.translated.update({string: translation})
                self.cloud_requests += 1

            else:
                self.translated.update({string: "%s(%s)" % (self.target_lang, string)})
                self.dummy_text += 1

            string = u'{}'.format(self.translated[string])

        elif string:
            string = u'{}'.format(self.translated[string])
            self.dict_hits += 1

        else:
            self.empty_strings += 1

        return string

    def multi_line(self, string):
        lines = string.split('\n')
        return "\n".join([self.translate(line) for line in lines])

    def save_history(self):
        """ Save the translation dictionary as a history file. """
        if self.history:
            try:
                with open(self.history_file, 'w') as f:
                    json.dump(self.translated, f)
            except (IOError, OSError):
                fallback_history_file = os.path.join(os.environ['HOME'], 'history_file.json')
                warnings.warn("'%s' is not writeable - saving history in home directory as:\n '%s' " % (
                    self.history_file, fallback_history_file
                ))
                self.history_file = fallback_history_file
                with open(self.history_file, 'w') as f:
                    json.dump(self.translated, f)

    def update_stats(self):
        """ Provide a set of summary stats on how the translation was done. """
        language_pair = '%s-%s' % (self.source_lang, self.target_lang)
        for stat in ['cloud_requests', 'dummy_text', 'dict_hits', 'empty_strings']:
            self.stats[language_pair].update({stat: getattr(self, stat)})
        self.stats[language_pair].update({'history': len(self.translated)})


class TranslateBase(object):
    """ Build translation framework, independent of file format. """
    def __init__(self, filepath, filename, translator, target=None, condense=False, cross_check=False):
        self.filepath = filepath
        self.source = filename
        self.filename, self.ext = filename.rsplit('.', 1)
        self.translator = translator
        self.condense = condense
        self.cross_check = hasattr(self.translator, 'source_lang') and bool(self.translator.source_lang) and cross_check
        if cross_check and not bool(self.translator.source_lang):
            warnings.warn('Not possible to translate back: no source language was given')

        if not target:
            self.target = self.add_lang_to_filename(self.source)
        else:
            self.target = target

    def execute(self, translate_method):
        """ Method to combine translation and cross_check as in subclass.__init__ for .docx etc """
        translate_method()
        self.translator.save_history()
        self.translator.update_stats()

        if self.cross_check:
            self.swap_languages()
            translate_method()
            self.translator.save_history()
            self.translator.update_stats()

        # Google does not close the session: Connection='keep-alive'.  The line below did not work.
        # self.translator.client._connection.http.close()

    def swap_languages(self):
        """ Swap languages and re-translate as a crude way of assessing the quality of the translation. """
        self.translator.target_lang, self.translator.source_lang = (
            self.translator.source_lang, self.translator.target_lang
        )
        if self.translator.history:
            self.translator.history_file = self.add_lang_to_filename(self.translator.history_file)
        self.source, self.target = (self.target, self.add_lang_to_filename(self.target))
        self.translator.prepare_translation()

    def add_lang_to_filename(self, filename, lang=None):
        """ Insert language code into a filename. """
        if not lang and hasattr(self.translator, 'target_lang'):
            lang = self.translator.target_lang

        if '.' in filename:
            body, extn = filename.rsplit('.', 1)
        else:
            body, extn = filename, self.ext

        return '%s_%s.%s' % (body, lang, extn)

    def set_language(self, document_object):
        """ Configure the language settings in the output document. """
        if self.translator.target_lang in LANGUAGE_PROPERTIES:
            if 'lang' in LANGUAGE_PROPERTIES[self.translator.target_lang]:
                document_object.core_properties.language = LANGUAGE_PROPERTIES[self.translator.target_lang]['lang']

            if not isinstance(self, TranslatePptx) and 'font' in LANGUAGE_PROPERTIES[self.translator.target_lang]:
                document_object.styles["Normal"].font.name = LANGUAGE_PROPERTIES[self.translator.target_lang]['font']

        # TODO: [future] self.document.styles["Normal"] language = LANGUAGE_PROPERTIES[self.translator.target_lang]['lang']
        # (if/when supported by python-docx)

    def condense_runs(self, paragraph, brk_run=None):
        if len(paragraph.runs) > 1:
            self._previous_run = paragraph.runs[0]
            for i, run in enumerate(paragraph.runs[1:]):
                if all((
                        self._previous_run.text,
                        run.text,
                        (not isinstance(self, TranslatePptx) or not brk_run[i + 1]),
                        self.same_style_runs(self._previous_run, run, paragraph, DOCX_STYLE_PROPERTY)
                )):
                    self._previous_run.text += run.text
                    run.text = ''
                else:
                    self._previous_run = run

    def translate_paragraphs(self, document_object):
        """ Perform the translation of each text element """
        for paragraph in document_object:
            translate_style = any((
                                      isinstance(self, TranslatePptx),
                                      isinstance(self, TranslateDocx)
                                      and paragraph.style.name not in TRANSLATION_RULES['ignore_styles'],
            ))
            brk_run = self.break_runs(paragraph) if isinstance(self, TranslatePptx) else None
            if translate_style:
                if self.condense:
                    self.condense_runs(paragraph, brk_run=brk_run)
                for run in paragraph.runs:
                    if run.text:  # and '\n' in run.text:
                        run.text = self.translator.multi_line(run.text)

    def same_style_runs(self, reference, comparison, para_style, attrs):
        """
        A method for comparing styles.

        Does not pick up any of the following:
        a. <w:proofErr w:type="spellStart"/> (desirable as a separate submission to avoid mis-translation)
        d. any differences other than the style attributes in DOCX_STYLE_PROPERTY
           (list based on @property in docx Run class)
        """
        for attr in attrs:
            if reference and comparison and hasattr(reference, attr) and hasattr(comparison, attr):
                sub_ref = getattr(reference, attr)
                sub_comp = getattr(comparison, attr)
                # if None, then need to look at the parent style
                if hasattr(para_style, attr):
                    par_comp = getattr(para_style, attr)
                    if sub_ref is None:
                        sub_ref = par_comp
                    if sub_comp is None:
                        sub_comp = par_comp
                else:
                    par_comp = None

                if attrs[attr]:
                    if not self.same_style_runs(sub_ref, sub_comp, par_comp, attrs[attr]):
                        return False
                elif sub_ref != sub_comp:
                    return False
                else:
                    pass

            elif reference and comparison and hasattr(reference, attr) != hasattr(comparison, attr):
                return False

            else:
                pass

        return True


class PreserveWhitespace(object):
    # TODO: incorporate in other translators
    def __init__(self, string):
        body = re.escape(string.strip())
        m = re.search(f'(?P<body>{body})', string)
        self.body = m["body"] if m else ''
        start, finish = m.span()
        self.leading_ws = string[:start]
        self.trailing_ws = string[finish:]
        self.full_text = self.body.join([self.leading_ws, self.trailing_ws])

    def replace(self, body):
        return body.join([self.leading_ws, self.trailing_ws])


class TranslateExcel(TranslateBase):
    """
    Translate text in an Excel (.xlsx) spreadsheet file

    Previously released version of openpyxl does not support preservation of embedded images.
    """
    # TODO: [future] check not removing embedded images, once this is available in openpyxl
    # TODO: [future] preserve rich text formatting within a cell if/when this is available in openpyxl
    # Note: MergeCell formatting is not working in openpyxl 2.5.5 but does work in development code
    def __init__(self, filepath, filename, translator, target=None, condense=False, cross_check=False):
        super(TranslateExcel, self).__init__(filepath, filename, translator,
                                             target=target, condense=condense, cross_check=cross_check)
        self.wb = load_workbook(os.path.join(self.filepath, self.source))
        self.execute(self.translate)

    def translate(self):
        """ Perform the translation of each text element """
        for sheetname in self.wb.sheetnames:
            for row in self.wb[sheetname]:
                for cell in row:
                    if isinstance(cell, Cell):
                        if cell.value and cell.data_type == 's':
                            cell.value = self.translator.multi_line(cell.value)

        self.wb.save(os.path.join(self.filepath, self.target))


class TranslateDocx(TranslateBase):
    """ Translate text in a Word (.docx) document file """
    def __init__(self, filepath, filename, translator, target=None, condense=False, cross_check=False):
        super(TranslateDocx, self).__init__(filepath, filename, translator,
                                            target=target, condense=condense, cross_check=cross_check)
        self.document = Document(os.path.join(self.filepath, self.source))
        self.execute(self.translate)

    def translate(self):
        """ Break the document up into text elements for translation """
        self.set_language(self.document)
        self.translate_paragraphs(self.document.paragraphs)
        for table in self.document.tables:
            for row in table.rows:
                for cell in row.cells:
                    self.translate_paragraphs(cell.paragraphs)

        # TODO: [future] translate headers/footers/text_boxes (if/when supported by python-docx)
        # for section in document.sections:
        #     # print(document.sections[section].footer)
        #     print(section._sectPr)

        self.document.save(os.path.join(self.filepath, self.target))


class TranslatePptx(TranslateBase):
    """ Translate text in a PowerPoint (.pptx) presentation file """
    def __init__(self, filepath, filename, translator, target=None, condense=False, cross_check=False):
        super(TranslatePptx, self).__init__(filepath, filename, translator,
                                            target=target, condense=condense, cross_check=cross_check)
        self.prs = Presentation(os.path.join(self.filepath, self.source))
        self.execute(self.translate)

    def translate(self):
        """ Perform the translation of each text element """
        self.set_language(self.prs)

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self.translate_paragraphs(shape.text_frame.paragraphs)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if hasattr(cell, 'text_frame'):
                                self.translate_paragraphs(cell.text_frame.paragraphs)

                # TODO: [future] add SmartArt text translation if/when supported by python-pptx

        self.prs.save(os.path.join(self.filepath, self.target))

    def break_runs(self, paragraph):
        if len(paragraph.runs) > 1:
            breaks = [isinstance(each, CT_TextLineBreak)
                      for each in paragraph._element.content_children]
            return [br for br in self.break_at_run(breaks)]

    def break_at_run(self, break_list):
        br = False
        for element in break_list:
            if element:
                br = True
            else:
                yield br
                br = False


class TranslateHtml(TranslateBase):
    # TODO: change lang attr e.g. <html class="no-js" lang="en-US">
    #  how does this appear in multi-language web sites so the browser can auto-select?
    """ Translate text in an HMTL (.html) file """
    def __init__(self, filepath, filename, translator, target=None, condense=False, cross_check=False):
        super(TranslateHtml, self).__init__(filepath, filename, translator,
                                            target=target, condense=condense, cross_check=cross_check)
        self.web_page = lxml.html.parse(os.path.join(self.filepath, self.source)).getroot()
        self.execute(self.translate)

    def translate(self):
        """ Translate method for html files to translate each element with text. """
        for attr in ['text', 'tail']:
            text = [
                element for element in self.web_page.iter()
                if not isinstance(element, lxml.html.HtmlComment)
                if element.tag not in ['html', 'head', 'meta', 'style', 'script']  #, 'a']
                if getattr(element, attr) is not None
                if getattr(element, attr).strip()
            ]
            for element in text:
                trimmed_text = PreserveWhitespace(getattr(element, attr))
                setattr(element, attr, trimmed_text.replace(self.translator.translate(trimmed_text.body)))
        with open(os.path.join(self.filepath, self.target), 'wb') as f:
            f.write(lxml.etree.tostring(self.web_page, method='html'))
            pass


class TranslateText(TranslateBase):
    """ Translate text in a plain text (e.g. .txt) file """
    def __init__(self, filepath, filename, translator, target=None, condense=False, cross_check=False):
        super(TranslateText, self).__init__(filepath, filename, translator,
                                            target=target, condense=condense, cross_check=cross_check)
        with open(os.path.join(self.filepath, self.source), 'r') as f:
            self.text = [line for line in f]
        self.execute(self.translate)

    def translate(self):
        """ Translate method for text files to translate each line of text. """
        with open(os.path.join(self.filepath, self.target), 'w') as f:
            for line in self.text:
                trimmed_text = PreserveWhitespace(line)
                f.write(trimmed_text.replace(self.translator.translate(trimmed_text.body)))


